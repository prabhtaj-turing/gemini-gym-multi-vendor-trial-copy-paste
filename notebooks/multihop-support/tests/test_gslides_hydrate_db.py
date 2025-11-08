import unittest
import os
import json
import tempfile
import shutil
from pathlib import Path
import random
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# Import the hydrate_db function from the gdrive utils
import sys
sys.path.append(os.path.join(os.path.dirname(__file__), '..','..', '..', 'APIs'))
from gdrive.SimulationEngine.utils import hydrate_db
from gdrive.SimulationEngine.db import DB

def hydrate_db_sample(db_instance, directory_path, max_files_per_type=3):
    """
    Load a sample of JSON files from the directory for testing purposes.
    
    Args:
        db_instance: The database instance to hydrate
        directory_path: Path to the test data directory
        max_files_per_type: Maximum number of files to load per file type
    
    Returns:
        bool: True if hydration was successful
    """
    if not os.path.isdir(directory_path):
        raise FileNotFoundError(f"Directory not found: '{directory_path}'")
    
    # Group files by type
    file_types = {
        'docs': [],
        'sheets': [],
        'slides': [],
        'other': []
    }
    
    # Walk through directory and categorize files
    for root, _, files in os.walk(directory_path):
        for file in files:
            if file.endswith('.json'):
                file_path = os.path.join(root, file)
                
                # Categorize by file type based on the filename
                if '.docx.json' in file:
                    file_types['docs'].append(file_path)
                elif '.xlsx.json' in file:
                    file_types['sheets'].append(file_path)
                elif '.pptx.json' in file:
                    file_types['slides'].append(file_path)
                else:
                    file_types['other'].append(file_path)
    
    # Sample files from each type
    sampled_files = []
    for file_type, files in file_types.items():
        if files:
            # Randomly sample up to max_files_per_type from each category
            sample_size = min(max_files_per_type, len(files))
            sampled_files.extend(random.sample(files, sample_size))
    
    # Load sampled files
    user_id = 'me'
    from gdrive.SimulationEngine.utils import _ensure_user
    _ensure_user(user_id)
    
    db_user = db_instance['users'][user_id]
    all_json_data = []
    
    for file_path in sampled_files:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                all_json_data.append(data)
        except json.JSONDecodeError:
            print(f"Warning: Could not decode JSON from file: {file_path}")
        except Exception as e:
            print(f"An error occurred while reading {file_path}: {e}")
    
    db_user['files'] = {file['id']: file for file in all_json_data}
    
    print(f"✓ Loaded {len(all_json_data)} sample files for testing")
    return True

class TestGSlidesHydrateDB(unittest.TestCase):
    """Tests for Google Slides hydrate_db functionality."""

    def setUp(self):
        """Reset DB to a clean state before each test."""
        # Clear the DB and set up a fresh state
        DB.clear()
        DB.update({
            'users': {
                'me': {
                    'about': {
                        'kind': 'drive#about',
                        'storageQuota': {
                            'limit': '1073741824',  # 1GB in bytes
                            'usageInDrive': '0',
                            'usageInDriveTrash': '0',
                            'usage': '0'
                        },
                        'driveThemes': False,
                        'canCreateDrives': False,
                        'importFormats': {},
                        'exportFormats': {},
                        'appInstalled': False,
                        'user': {
                            'displayName': 'Test User',
                            'kind': 'drive#user',
                            'me': True,
                            'permissionId': 'test_perm_id',
                            'emailAddress': 'test@example.com'
                        },
                        'folderColorPalette': "",
                        'maxImportSizes': {},
                        'maxUploadSize': '104857600'  # 100MB in bytes
                    },
                    'files': {},
                    'drives': {},
                    'comments': {},
                    'replies': {},
                    'labels': {},
                    'accessproposals': {},
                    'counters': {
                        'file': 0,
                        'drive': 0,
                        'comment': 0,
                        'reply': 0,
                        'label': 0,
                        'accessproposal': 0,
                        'revision': 0
                    }
                }
            }
        })

    def get_test_data_path(self):
        """Get the path to the test data directory."""
        return os.path.join(os.path.dirname(__file__), 'data')

    def test_google_slides_content_integrity(self):
        """Test that Google Slides files are properly loaded with content."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find Google Slides files
        slides_files = []
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_files.append((file_id, file_data))
        
        # Should have at least one slides file
        self.assertGreater(len(slides_files), 0, "Should have at least one Google Slides file")
        
        # Test each slides file
        for file_id, file_data in slides_files:
            with self.subTest(file_id=file_id):
                # Verify required fields
                self.assertIn('id', file_data)
                self.assertIn('name', file_data)
                self.assertIn(file_data['mimeType'], ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'])
                
                # Verify slides structure
                if 'slides' in file_data:
                    self.assertIsInstance(file_data['slides'], list)
                    
                # Verify presentation ID
                if 'presentationId' in file_data:
                    self.assertEqual(file_data['presentationId'], file_data['id'])
                    
                print(f"✓ Validated Google Slides file: {file_data['name']}")

    def test_slides_structure_validation(self):
        """Test the structure of Google Slides presentations."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find Google Slides files
        slides_files = []
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_files.append((file_id, file_data))
        
        if not slides_files:
            self.fail("No Google Slides files found in test data")
        
        for file_id, file_data in slides_files:
            with self.subTest(file_id=file_id):
                # Test slides structure
                if 'slides' in file_data:
                    for slide in file_data['slides']:
                        self.assertIn('objectId', slide)
                        self.assertIn('pageType', slide)
                        
                        # Verify pageType is valid
                        if 'pageType' in slide:
                            valid_page_types = ['SLIDE', 'MASTER', 'LAYOUT']
                            self.assertIn(slide['pageType'], valid_page_types)
                            
                # Test presentation metadata
                if 'title' in file_data:
                    self.assertIsInstance(file_data['title'], str)
                    
                if 'pageSize' in file_data:
                    self.assertIsInstance(file_data['pageSize'], dict)
                    
                print(f"✓ Validated Google Slides structure: {file_data['name']}")

    def test_slides_page_elements(self):
        """Test that Google Slides page elements are properly loaded."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find Google Slides files
        slides_files = []
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_files.append((file_id, file_data))
        
        if not slides_files:
            self.fail("No Google Slides files found in test data")
        
        for file_id, file_data in slides_files:
            with self.subTest(file_id=file_id):
                if 'slides' in file_data:
                    for slide in file_data['slides']:
                        # Check for page elements
                        if 'pageElements' in slide:
                            self.assertIsInstance(slide['pageElements'], list)
                            
                            for element in slide['pageElements']:
                                self.assertIn('objectId', element)
                                
                print(f"✓ Validated Google Slides page elements: {file_data['name']}")

    def test_slides_presentation_properties(self):
        """Test that Google Slides presentation properties are correctly preserved."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find Google Slides files
        slides_files = []
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_files.append((file_id, file_data))
        
        if not slides_files:
            self.fail("No Google Slides files found in test data")
        
        for file_id, file_data in slides_files:
            with self.subTest(file_id=file_id):
                # Verify presentation properties
                if 'pageSize' in file_data:
                    page_size = file_data['pageSize']
                    if 'width' in page_size:
                        self.assertIsInstance(page_size['width'], dict)
                    if 'height' in page_size:
                        self.assertIsInstance(page_size['height'], dict)
                
                # Verify masters and layouts if present
                if 'masters' in file_data:
                    self.assertIsInstance(file_data['masters'], list)
                    
                if 'layouts' in file_data:
                    self.assertIsInstance(file_data['layouts'], list)
                
                print(f"✓ Validated Google Slides properties: {file_data['name']}")

    def test_api_integration_google_slides(self):
        """Test that hydrated Google Slides files can be accessed via the API."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find a Google Slides file
        slides_file_id = None
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_file_id = file_id
                break
        
        if slides_file_id is None:
            self.fail("No Google Slides file found in test data")
        
        # Test API access
        try:
            from google_slides import get_presentation
            presentation = get_presentation(presentationId=slides_file_id)
            
            # Verify API response
            self.assertIsInstance(presentation, dict)
            self.assertIn('presentationId', presentation)
            self.assertEqual(presentation['presentationId'], slides_file_id)
            
            print(f"✓ API integration test passed for Google Slides file: {slides_file_id}")
            
        except ImportError:
            self.fail("Google Slides API not available")
        except Exception as e:
            self.fail(f"API integration test failed: {e}")

    def test_slides_count_and_statistics(self):
        """Test the count and statistics of Google Slides files."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find Google Slides files
        slides_files = []
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_files.append((file_id, file_data))
        
        # Test that we have a reasonable number of slides files
        self.assertGreaterEqual(len(slides_files), 1, "Should have at least one Google Slides file")
        self.assertLessEqual(len(slides_files), 2, "Should not have more than 2 Google Slides files (sample limit)")
        
        # Test each slides file
        for file_id, file_data in slides_files:
            with self.subTest(file_id=file_id):
                # Verify basic structure
                self.assertIn('id', file_data)
                self.assertIn('name', file_data)
                self.assertIn(file_data['mimeType'], ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'])
                
                print(f"✓ Validated Google Slides statistics: {file_data['name']}")

    def test_slides_text_content(self):
        """Test that Google Slides text content is properly loaded."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find Google Slides files
        slides_files = []
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_files.append((file_id, file_data))
        
        if not slides_files:
            self.fail("No Google Slides files found in test data")
        
        for file_id, file_data in slides_files:
            with self.subTest(file_id=file_id):
                # Test text content in slides
                if 'slides' in file_data:
                    for slide in file_data['slides']:
                        if 'pageElements' in slide:
                            for element in slide['pageElements']:
                                if 'shape' in element and 'text' in element['shape']:
                                    self.assertIn('textElements', element['shape']['text'])
                                    
                print(f"✓ Validated Google Slides text content: {file_data['name']}")

    def test_slides_file_size_validation(self):
        """Test that Google Slides file sizes are reasonable."""
        test_data_path = self.get_test_data_path()
        
        if not os.path.exists(test_data_path):
            self.fail(f"Test data directory not found: {test_data_path}")
        
        # Hydrate the database with sample data
        hydrate_db_sample(DB, test_data_path, max_files_per_type=2)
        
        # Find Google Slides files
        slides_files = []
        for file_id, file_data in DB['users']['me']['files'].items():
            if file_data.get('mimeType') in ['application/vnd.google-apps.presentation', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                slides_files.append((file_id, file_data))
        
        if not slides_files:
            self.fail("No Google Slides files found in test data")
        
        for file_id, file_data in slides_files:
            with self.subTest(file_id=file_id):
                # Verify file size if present
                if 'size' in file_data:
                    size = int(file_data['size'])
                    self.assertGreater(size, 0, "File size should be greater than 0")
                    self.assertLess(size, 100 * 1024 * 1024, "File size should be less than 100MB")
                
                print(f"✓ Validated Google Slides file size: {file_data['name']}")

    def test_extracts_all_text_types(self):
        """Test that all text (paragraphs, runs, tables, groups) is extracted from a slide."""
        # Create a presentation in memory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Add a text box with two paragraphs and two runs
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = txBox.text_frame
        p1 = tf.add_paragraph()
        p1.text = "Paragraph 1"
        p2 = tf.add_paragraph()
        run1 = p2.add_run()
        run1.text = "Run 1"
        run2 = p2.add_run()
        run2.text = "Run 2"
        # Add a table
        table = slide.shapes.add_table(1, 2, Inches(1), Inches(2.5), Inches(2), Inches(0.5)).table
        table.cell(0, 0).text = "Cell 1"
        table.cell(0, 1).text = "Cell 2"
        # Save to a BytesIO buffer
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        # Write to a temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(buf.read())
            tmp_path = tmp.name
        # Import the converter (fix sys.path for parent dir)
        import sys, os
        parent_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
        if parent_dir not in sys.path:
            sys.path.insert(0, parent_dir)
        from gslides_converter import convert_pptx_to_gslides_format
        result = convert_pptx_to_gslides_format(tmp_path)
        os.unlink(tmp_path)
        # Gather all text from textElements
        all_text = []
        for slide in result['slides']:
            for el in slide['pageElements']:
                if 'shape' in el and 'text' in el['shape']:
                    for te in el['shape']['text'].get('textElements', []):
                        if 'textRun' in te:
                            all_text.append(te['textRun']['content'])
        
        # Check for all expected text
        self.assertIn("Paragraph 1", all_text)
        self.assertIn("Run 1", all_text)
        self.assertIn("Run 2", all_text)
        self.assertIn("Cell 1", all_text)
        self.assertIn("Cell 2", all_text)

    def test_real_pptx_text_content_matches_reference(self):
        """Test that all key text content from the reference JSON is present in the converter output."""
        import sys, os, json
        parent_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
        if parent_dir not in sys.path:
            sys.path.insert(0, parent_dir)
        from gslides_converter import convert_pptx_to_gslides_format
        pptx_path = os.path.join(parent_dir, 'input-data', 'files',
            'MP-42825 Files-20250619T235923Z-1-001', 'MP-42825 Files', 'F125', 'The Creative Canvas',
            '2025 Summer Season_ Summer Spark!', 'Summer_Spark!_2025_seasonal_overview.pptx')
        ref_json_path = os.path.join(parent_dir, 'output-data', 'files',
            'MP-42825 Files-20250619T235923Z-1-001', 'MP-42825 Files', 'F125', 'The Creative Canvas',
            '2025 Summer Season_ Summer Spark!', 'Summer_Spark!_2025_seasonal_overview.pptx.json')
        # Convert PPTX
        result = convert_pptx_to_gslides_format(pptx_path)
        # Load reference JSON
        with open(ref_json_path, 'r', encoding='utf-8') as f:
            ref = json.load(f)
        # Helper to extract and normalize all text from slides
        def extract_and_normalize(slides):
            texts = []
            for slide in slides:
                for el in slide.get('pageElements', []):
                    shape = el.get('shape', {})
                    text = shape.get('text', {})
                    for te in text.get('textElements', []):
                        if 'textRun' in te:
                            # Normalize: strip whitespace and newlines
                            t = te['textRun']['content'].replace('\n', ' ').strip()
                            if t:
                                texts.append(t)
            # Join all text runs into a single string for comparison
            return ' '.join(texts)
        result_text = extract_and_normalize(result['slides'])
        ref_text = extract_and_normalize(ref['slides'])
        # Extract key phrases from reference text (non-empty, meaningful content)
        ref_phrases = [phrase.strip() for phrase in ref_text.split() if len(phrase.strip()) > 3]
        # Check that all key phrases from reference are present in result
        missing_phrases = []
        for phrase in ref_phrases:
            # Normalize phrase by removing punctuation for comparison
            normalized_phrase = phrase.rstrip('.,!?;:')
            if normalized_phrase not in result_text:
                missing_phrases.append(phrase)
        # Assert that all key reference content is present in the result
        self.assertEqual(len(missing_phrases), 0, 
            f"Missing key phrases from reference in converter output: {missing_phrases}\n"
            f"Reference text: {ref_text}\n"
            f"Result text: {result_text[:500]}...")
        # Also verify that the result contains more content (enhanced extraction)
        self.assertGreater(len(result_text), len(ref_text), 
            f"Converter should extract more content than reference. "
            f"Reference length: {len(ref_text)}, Result length: {len(result_text)}")

if __name__ == '__main__':
    unittest.main() 
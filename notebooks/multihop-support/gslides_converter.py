import os
import uuid
import hashlib
from datetime import datetime, UTC
from pptx import Presentation


def convert_pptx_to_gslides_format(file_path, file_name=None):
    """Convert a PowerPoint file to Google Slides format JSON."""
    if file_name is None:
        file_name = os.path.basename(file_path)
    
    file_stats = os.stat(file_path)
    presentation_id = f"pres_{hashlib.md5(file_path.encode()).hexdigest()}"
    current_time = datetime.fromtimestamp(file_stats.st_mtime, UTC).strftime('%Y-%m-%dT%H:%M:%SZ')
    
    try:
        prs = Presentation(file_path)
        slides = []
        
        # Extract presentation metadata first
        metadata_elements = _extract_presentation_metadata(prs)
        
        for i, slide in enumerate(prs.slides):
            page_elements = []
            
            # Add speaker notes as additional text elements
            speaker_notes = _extract_speaker_notes(slide)
            
            for j, shape in enumerate(slide.shapes):
                element = {
                    "objectId": f"element{j+1}_slide{i+1}",
                    "size": {
                        "width": {
                            "magnitude": float(getattr(shape, "width", 200)),
                            "unit": "PT"
                        },
                        "height": {
                            "magnitude": float(getattr(shape, "height", 100)),
                            "unit": "PT"
                        }
                    },
                    "transform": {
                        "scaleX": 1.0,
                        "scaleY": 1.0,
                        "translateX": float(getattr(shape, "left", 0)),
                        "translateY": float(getattr(shape, "top", 0)),
                        "unit": "PT"
                    },
                    "shape": {
                        "shapeType": _map_shape_type(getattr(shape, "shape_type", None)),
                        "text": {
                            "textElements": _extract_text_elements(shape)
                        }
                    }
                }
                page_elements.append(element)
            
            # Add speaker notes as a separate element if they exist
            if speaker_notes:
                notes_element = {
                    "objectId": f"notes_element_slide{i+1}",
                    "size": {
                        "width": {"magnitude": 600, "unit": "PT"},
                        "height": {"magnitude": 100, "unit": "PT"}
                    },
                    "transform": {
                        "scaleX": 1.0,
                        "scaleY": 1.0,
                        "translateX": 50,
                        "translateY": 500,
                        "unit": "PT"
                    },
                    "shape": {
                        "shapeType": "TEXT_BOX",
                        "text": {
                            "textElements": speaker_notes
                        }
                    }
                }
                page_elements.append(notes_element)
            
            slides.append({
                "objectId": f"slide{i+1}_page{i+1}",
                "pageType": "SLIDE",
                "pageProperties": {
                    "backgroundColor": {
                        "opaqueColor": {
                            "rgbColor": {
                                "red": 1.0,
                                "green": 1.0,
                                "blue": 1.0
                            }
                        }
                    }
                },
                "slideProperties": {
                    "masterObjectId": "master1",
                    "layoutObjectId": "layout1"
                },
                "pageElements": page_elements,
                "revisionId": f"rev_slide{i+1}"
            })
        
        # Add a metadata slide at the beginning if metadata exists
        if metadata_elements:
            metadata_slide = {
                "objectId": "metadata_slide",
                "pageType": "SLIDE",
                "pageProperties": {
                    "backgroundColor": {
                        "opaqueColor": {
                            "rgbColor": {
                                "red": 0.95,
                                "green": 0.95,
                                "blue": 0.95
                            }
                        }
                    }
                },
                "slideProperties": {
                    "masterObjectId": "master1",
                    "layoutObjectId": "layout1"
                },
                "pageElements": [{
                    "objectId": "metadata_element",
                    "size": {
                        "width": {"magnitude": 600, "unit": "PT"},
                        "height": {"magnitude": 400, "unit": "PT"}
                    },
                    "transform": {
                        "scaleX": 1.0,
                        "scaleY": 1.0,
                        "translateX": 50,
                        "translateY": 50,
                        "unit": "PT"
                    },
                    "shape": {
                        "shapeType": "TEXT_BOX",
                        "text": {
                            "textElements": metadata_elements
                        }
                    }
                }],
                "revisionId": "rev_metadata"
            }
            slides.insert(0, metadata_slide)
        
        # Create the Drive file structure with Google Slides data embedded
        # This matches GoogleSlidesDefaultDB.json format exactly
        json_data = {
            "id": presentation_id,
            "driveId": "",
            "name": file_name,
            "mimeType": "application/vnd.google-apps.presentation",
            "createdTime": datetime.fromtimestamp(file_stats.st_ctime, UTC).strftime('%Y-%m-%dT%H:%M:%SZ'),
            "modifiedTime": current_time,
            "trashed": False,
            "starred": False,
            "parents": [],
            "owners": ["john.doe@gmail.com"],
            "size": str(file_stats.st_size),
            "permissions": [
                {
                    "id": f"permission_{presentation_id}",
                    "role": "owner",
                    "type": "user",
                    "emailAddress": "john.doe@gmail.com"
                }
            ],
            # Google Slides specific fields embedded in the Drive file
            "presentationId": presentation_id,
            "title": os.path.splitext(file_name)[0],
            "pageSize": {
                "width": {
                    "magnitude": float(prs.slide_width),
                    "unit": "EMU"
                },
                "height": {
                    "magnitude": float(prs.slide_height),
                    "unit": "EMU"
                }
            } if hasattr(prs, 'slide_width') and hasattr(prs, 'slide_height') else None,
            "slides": slides,
            "masters": [],
            "layouts": [],
            "notesMaster": None,
            "locale": "en-US",
            "revisionId": f"rev_{presentation_id}"
        }
        
    except Exception as e:
        # Create error structure that still matches the expected schema
        json_data = {
            "id": presentation_id,
            "driveId": "",
            "name": file_name,
            "mimeType": "application/vnd.google-apps.presentation",
            "createdTime": datetime.fromtimestamp(file_stats.st_ctime, UTC).strftime('%Y-%m-%dT%H:%M:%SZ'),
            "modifiedTime": current_time,
            "trashed": False,
            "starred": False,
            "parents": [],
            "owners": ["john.doe@gmail.com"],
            "size": str(file_stats.st_size),
            "permissions": [
                {
                    "id": f"permission_{presentation_id}",
                    "role": "owner",
                    "type": "user",
                    "emailAddress": "john.doe@gmail.com"
                }
            ],
            "presentationId": presentation_id,
            "title": os.path.splitext(file_name)[0],
            "pageSize": None,
            "slides": [],
            "masters": [],
            "layouts": [],
            "notesMaster": None,
            "locale": "en-US",
            "revisionId": f"rev_{presentation_id}",
            "error": f"Error reading PowerPoint file: {e}"
        }
    
    return json_data


def _map_shape_type(pptx_shape_type):
    """Map python-pptx shape types to Google Slides shape types."""
    # Map common shape types - this is a simplified mapping
    shape_type_map = {
        1: "TEXT_BOX",      # MSO_SHAPE_TYPE.TEXT_BOX
        5: "RECTANGLE",     # MSO_SHAPE_TYPE.AUTO_SHAPE rectangle
        9: "ELLIPSE",       # MSO_SHAPE_TYPE.AUTO_SHAPE oval/ellipse
        # Add more mappings as needed
    }
    
    if pptx_shape_type is None:
        return "TEXT_BOX"
    
    return shape_type_map.get(pptx_shape_type, "TEXT_BOX")


def _extract_text_elements(shape):
    """Recursively extract all text elements from a shape, including paragraphs, runs, tables, and groups."""
    text_elements = []
    
    # Text frames (text boxes, placeholders, etc.)
    if hasattr(shape, 'text_frame') and shape.text_frame is not None:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text_elements.append({
                    "textRun": {
                        "content": run.text,
                        "style": {
                            "fontFamily": run.font.name if run.font and run.font.name else "Calibri",
                            "fontSize": {
                                "magnitude": run.font.size.pt if run.font and run.font.size else 12,
                                "unit": "PT"
                            }
                        }
                    }
                })
            # If paragraph has no runs but has text, add as a single run
            if not para.runs and para.text:
                text_elements.append({
                    "textRun": {
                        "content": para.text,
                        "style": {
                            "fontFamily": "Calibri",
                            "fontSize": {"magnitude": 12, "unit": "PT"}
                        }
                    }
                })
    
    # Tables
    if hasattr(shape, 'table') and shape.table is not None:
        for row in shape.table.rows:
            for cell in row.cells:
                # Recursively extract from each cell (cells are shapes)
                text_elements.extend(_extract_text_elements(cell))
    
    # Grouped shapes
    if hasattr(shape, 'shapes'):
        for subshape in shape.shapes:
            text_elements.extend(_extract_text_elements(subshape))
    
    # Extract alt text if available
    if hasattr(shape, 'alt_text') and shape.alt_text:
        text_elements.append({
            "textRun": {
                "content": f"[Alt text: {shape.alt_text}]",
                "style": {
                    "fontFamily": "Calibri",
                    "fontSize": {"magnitude": 10, "unit": "PT"}
                }
            }
        })
    
    # Extract chart data labels and titles if available
    try:
        if hasattr(shape, 'chart') and shape.chart is not None:
            chart = shape.chart
            # Chart title
            if hasattr(chart, 'chart_title') and chart.chart_title and hasattr(chart.chart_title, 'text_frame') and chart.chart_title.text_frame:
                for para in chart.chart_title.text_frame.paragraphs:
                    if para.text.strip():
                        text_elements.append({
                            "textRun": {
                                "content": f"[Chart title: {para.text}]",
                                "style": {
                                    "fontFamily": "Calibri",
                                    "fontSize": {"magnitude": 12, "unit": "PT"}
                                }
                            }
                        })
            
            # Axis titles
            if hasattr(chart, 'axis_title') and chart.axis_title and hasattr(chart.axis_title, 'text_frame') and chart.axis_title.text_frame:
                for para in chart.axis_title.text_frame.paragraphs:
                    if para.text.strip():
                        text_elements.append({
                            "textRun": {
                                "content": f"[Axis title: {para.text}]",
                                "style": {
                                    "fontFamily": "Calibri",
                                    "fontSize": {"magnitude": 10, "unit": "PT"}
                                }
                            }
                        })
    except Exception:
        # Ignore chart extraction errors
        pass
    
    # Extract SmartArt text if available
    try:
        if hasattr(shape, 'smart_art') and shape.smart_art is not None:
            # SmartArt diagrams can contain text in their shapes
            if hasattr(shape.smart_art, 'all_shapes'):
                for smart_shape in shape.smart_art.all_shapes:
                    if hasattr(smart_shape, 'text_frame') and smart_shape.text_frame:
                        for para in smart_shape.text_frame.paragraphs:
                            if para.text.strip():
                                text_elements.append({
                                    "textRun": {
                                        "content": f"[SmartArt: {para.text}]",
                                        "style": {
                                            "fontFamily": "Calibri",
                                            "fontSize": {"magnitude": 11, "unit": "PT"}
                                        }
                                    }
                                })
    except Exception:
        # Ignore SmartArt extraction errors
        pass
    
    return text_elements 


def _extract_speaker_notes(slide):
    """Extract speaker notes from a slide if available."""
    text_elements = []
    
    # Check if the slide has notes
    if hasattr(slide, 'notes_slide') and slide.notes_slide is not None:
        notes_slide = slide.notes_slide
        if hasattr(notes_slide, 'shapes'):
            for shape in notes_slide.shapes:
                # Extract text from notes shapes
                if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                text_elements.append({
                                    "textRun": {
                                        "content": f"[Speaker note: {run.text}]",
                                        "style": {
                                            "fontFamily": "Calibri",
                                            "fontSize": {"magnitude": 10, "unit": "PT"}
                                        }
                                    }
                                })
                        # If paragraph has no runs but has text
                        if not para.runs and para.text.strip():
                            text_elements.append({
                                "textRun": {
                                    "content": f"[Speaker note: {para.text}]",
                                    "style": {
                                        "fontFamily": "Calibri",
                                        "fontSize": {"magnitude": 10, "unit": "PT"}
                                    }
                                }
                            })
    
    return text_elements


def _extract_presentation_metadata(prs):
    """Extract metadata and other text content from the presentation."""
    text_elements = []
    
    # Extract presentation properties
    if hasattr(prs, 'core_properties'):
        core_props = prs.core_properties
        if core_props.title:
            text_elements.append({
                "textRun": {
                    "content": f"[Presentation title: {core_props.title}]",
                    "style": {
                        "fontFamily": "Calibri",
                        "fontSize": {"magnitude": 14, "unit": "PT"}
                    }
                }
            })
        if core_props.subject:
            text_elements.append({
                "textRun": {
                    "content": f"[Subject: {core_props.subject}]",
                    "style": {
                        "fontFamily": "Calibri",
                        "fontSize": {"magnitude": 12, "unit": "PT"}
                    }
                }
            })
        if core_props.author:
            text_elements.append({
                "textRun": {
                    "content": f"[Author: {core_props.author}]",
                    "style": {
                        "fontFamily": "Calibri",
                        "fontSize": {"magnitude": 10, "unit": "PT"}
                    }
                }
            })
        if core_props.comments:
            text_elements.append({
                "textRun": {
                    "content": f"[Comments: {core_props.comments}]",
                    "style": {
                        "fontFamily": "Calibri",
                        "fontSize": {"magnitude": 10, "unit": "PT"}
                    }
                }
            })
    
    # Extract slide master content (headers, footers, etc.)
    if hasattr(prs, 'slide_masters') and prs.slide_masters:
        for i, master in enumerate(prs.slide_masters):
            for shape in master.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                text_elements.append({
                                    "textRun": {
                                        "content": f"[Master {i+1}: {run.text}]",
                                        "style": {
                                            "fontFamily": "Calibri",
                                            "fontSize": {"magnitude": 9, "unit": "PT"}
                                        }
                                    }
                                })
    
    return text_elements 